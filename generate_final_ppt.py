from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

# ─── Color Palette ───────────────────────────────────────────────
BG_DARK     = RGBColor(10,  17,  40)   # navy dark
BG_CARD     = RGBColor(18,  32,  70)   # card bg
ACCENT_BLUE = RGBColor(56, 189, 248)   # sky-400
ACCENT_TEAL = RGBColor(45, 212, 191)   # teal-400
ACCENT_INDIGO = RGBColor(129, 140, 248) # indigo-400
WHITE       = RGBColor(248, 250, 252)
SLATE_300   = RGBColor(203, 213, 225)
SLATE_400   = RGBColor(148, 163, 184)

def set_bg(slide, color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def inch(x): return Inches(x)

# ─── Slide Helpers ────────────────────────────────────────────────

def slide_title(prs, title, sub, accent=ACCENT_BLUE):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    set_bg(slide, BG_DARK)

    # Accent bar left
    add_rect(slide, inch(0), inch(0), inch(0.12), inch(7.5), accent)

    # Title
    add_textbox(slide, title,
                inch(0.35), inch(2.6), inch(9.3), inch(1.2),
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Underline accent bar
    add_rect(slide, inch(0.35), inch(3.85), inch(4.5), inch(0.06), accent)

    # Sub
    add_textbox(slide, sub,
                inch(0.35), inch(3.95), inch(9.3), inch(1.0),
                font_size=22, color=SLATE_300, align=PP_ALIGN.LEFT)

    # Top badge
    add_textbox(slide, "🦈  Shark Task  |  Software Engineering",
                inch(0.35), inch(0.2), inch(9.3), inch(0.4),
                font_size=12, color=SLATE_400, align=PP_ALIGN.LEFT)
    return slide

def slide_section(prs, section_title, content_lines, accent=ACCENT_BLUE, icon=""):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_bg(slide, BG_DARK)

    # top accent bar
    add_rect(slide, inch(0), inch(0), inch(10), inch(0.09), accent)

    # section label
    add_textbox(slide, "🦈  Shark Task  |  Software Engineering",
                inch(0.4), inch(0.15), inch(9), inch(0.3),
                font_size=11, color=SLATE_400)

    # title
    add_textbox(slide, f"{icon}  {section_title}" if icon else section_title,
                inch(0.4), inch(0.55), inch(9.2), inch(0.9),
                font_size=34, bold=True, color=WHITE)

    # accent divider
    add_rect(slide, inch(0.4), inch(1.45), inch(1.5), inch(0.05), accent)

    # content cards
    card_top = inch(1.6)
    for line in content_lines:
        if line.startswith("##"):  # sub-heading
            add_textbox(slide, line[2:].strip(),
                        inch(0.4), card_top, inch(9.2), inch(0.5),
                        font_size=16, bold=True, color=accent)
            card_top += inch(0.45)
        elif line.startswith("---"):  # spacer
            card_top += inch(0.15)
        else:
            add_textbox(slide, line,
                        inch(0.55), card_top, inch(9), inch(0.42),
                        font_size=17, color=SLATE_300)
            card_top += inch(0.44)
    return slide

def slide_two_col(prs, title, left_lines, right_lines, accent=ACCENT_BLUE):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_bg(slide, BG_DARK)
    add_rect(slide, inch(0), inch(0), inch(10), inch(0.09), accent)
    add_textbox(slide, "🦈  Shark Task  |  Software Engineering",
                inch(0.4), inch(0.15), inch(9), inch(0.3),
                font_size=11, color=SLATE_400)
    add_textbox(slide, title,
                inch(0.4), inch(0.5), inch(9.2), inch(0.9),
                font_size=32, bold=True, color=WHITE)
    add_rect(slide, inch(0.4), inch(1.38), inch(1.5), inch(0.05), accent)

    # left card bg
    add_rect(slide, inch(0.3), inch(1.55), inch(4.45), inch(5.5),
             BG_CARD, accent)
    y = inch(1.75)
    for line in left_lines:
        if line.startswith("##"):
            add_textbox(slide, line[2:].strip(), inch(0.45), y, inch(4.1), inch(0.45),
                        font_size=14, bold=True, color=accent)
            y += inch(0.42)
        else:
            add_textbox(slide, line, inch(0.5), y, inch(4.0), inch(0.42),
                        font_size=15, color=SLATE_300)
            y += inch(0.4)

    # right card bg
    add_rect(slide, inch(5.0), inch(1.55), inch(4.6), inch(5.5),
             BG_CARD, ACCENT_TEAL)
    y = inch(1.75)
    for line in right_lines:
        if line.startswith("##"):
            add_textbox(slide, line[2:].strip(), inch(5.15), y, inch(4.3), inch(0.45),
                        font_size=14, bold=True, color=ACCENT_TEAL)
            y += inch(0.42)
        else:
            add_textbox(slide, line, inch(5.2), y, inch(4.2), inch(0.42),
                        font_size=15, color=SLATE_300)
            y += inch(0.4)
    return slide

# ─── Build Presentation ──────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 01  Title ──────────────────────────────────────────────────
    slide_title(prs,
        "Shark Task",
        "กระบวนการพัฒนาซอฟต์แวร์ตั้งแต่ต้นจนจบ\n"
        "Software Engineering Final Phase — Quality · Automation · Performance",
        ACCENT_BLUE)

    # ── 02  Agenda ─────────────────────────────────────────────────
    slide_section(prs, "สิ่งที่จะนำเสนอวันนี้", [
        "① กระบวนการทำงาน (SDLC & Workflow)",
        "② การออกแบบระบบ (Architecture)",
        "③ การทดสอบคุณภาพ (Testing — Live Demo)",
        "④ การวิเคราะห์ประสิทธิภาพ (Static & Dynamic Profiling)",
        "⑤ ระบบอัตโนมัติ (CI/CD Pipeline)",
        "⑥ Demo เว็บไซต์จริง",
    ], ACCENT_BLUE, "📋")

    # ── 03  SDLC Process ───────────────────────────────────────────
    slide_section(prs, "SDLC: กระบวนการพัฒนาตั้งแต่ต้น", [
        "## โมเดลที่ใช้:  Iterative & Incremental",
        "---",
        "🔵 Phase 1 — Requirements & Design",
        "   วิเคราะห์ความต้องการ, ออกแบบ DB Schema, วาง API Contract",
        "🟡 Phase 2 — Build (Core Features)",
        "   Auth → Project/Task CRUD → RBAC → Real-time Chat",
        "🟢 Phase 3 — Testing & Integration",
        "   Integration Tests 29 Cases · Coverage 83.6% · All PASS",
        "🔴 Phase 4 — Profiling & CI/CD",
        "   Static + Dynamic Profiling · GitHub Actions Pipeline",
    ], ACCENT_INDIGO, "⚙️")

    # ── 04  Requirement → Design Flow ─────────────────────────────
    slide_section(prs, "Phase 1: Requirements → Design", [
        "## ความต้องการหลัก (Requirements)",
        '▸  ระบบจัดการงานแบบ Role-Based (Manager / Worker / Admin)',
        '▸  ติดตาม Timeline แบบ Gantt Chart',
        '▸  ระบบแจ้งเตือนและ Activity Log',
        "---",
        "## การออกแบบที่ตามมา (Design Decisions)",
        '▸  PostgreSQL + Prisma  →  จัดการ Relational Data ที่ซับซ้อน',
        '▸  Next.js App Router  →  แยก Public / Dashboard Route ชัดเจน',
        '▸  JWT + Middleware    →  Security Layer แยกจาก Business Logic',
    ], ACCENT_TEAL, "📐")

    # ── 05  Build Phase ────────────────────────────────────────────
    slide_section(prs, "Phase 2: Build — ลำดับการสร้างระบบ", [
        "## Sprint 1  —  Foundation",
        "  DB Schema → Auth API (Register/Login/JWT) → Middleware",
        "## Sprint 2  —  Core Feature",
        "  Project CRUD → Task CRUD → RBAC → Task Status Machine",
        "## Sprint 3  —  Collaboration",
        "  Task Comments → File Attach → Activity Inbox → Messaging",
        "## Sprint 4  —  Polish",
        "  Analytics Dashboard → Gantt Chart → Admin Panel → Search/Filter",
        "---",
        "ทุก Sprint: เขียนโค้ด → ทดสอบ Manual → Commit → Review",
    ], ACCENT_BLUE, "🏗️")

    # ── 06  Architecture ───────────────────────────────────────────
    slide_two_col(prs,
        "สถาปัตยกรรมระบบ (System Architecture)",
        [
            "## Frontend",
            "▸ Next.js 14 App Router",
            "▸ Axios Interceptor (Auto JWT)",
            "▸ Tailwind CSS + CSS Variables",
            "▸ Socket.io Client (Real-time)",
            "---",
            "## ทำไม Next.js?",
            "▸ SSR ลด load time",
            "▸ API Route isolation",
            "▸ Middleware ดัก Auth ได้ก่อน render",
        ],
        [
            "## Backend",
            "▸ Node.js + Express",
            "▸ Prisma ORM → PostgreSQL",
            "▸ Socket.io Server",
            "▸ Multer (File Upload)",
            "---",
            "## หลักการ: Separation of Concerns",
            "▸ Controller แยกจาก DB Logic",
            "▸ Middleware ดูแล Auth แยกต่างหาก",
            "▸ Mock-able = Test-able",
        ],
        ACCENT_BLUE)

    # ── 07  Testing Strategy ───────────────────────────────────────
    slide_section(prs, "Phase 3: Testing Strategy", [
        "## แนวทาง: Black-box Integration Testing",
        "▸  ทดสอบพฤติกรรม API จากภายนอก เหมือนผู้ใช้จริง",
        "▸  ไม่ยึดติดกับ Implementation ภายใน (ทนต่อ Refactor)",
        "---",
        "## เครื่องมือ",
        "▸  Jest          — Test Runner + Assertion",
        "▸  Supertest     — จำลอง HTTP Request โดยไม่ต้องเปิด Server จริง",
        "▸  Manual Mock   — แทน Prisma Client ด้วย Mock เพื่อความเร็วและความแน่นอน",
        "---",
        "▸  ผลลัพธ์: 29 Tests · 3 Suites · ใช้เวลา ~1 วินาที · PASS 100%",
    ], ACCENT_TEAL, "🧪")

    # ── 08  Test Cases ─────────────────────────────────────────────
    slide_section(prs, "สิ่งที่ทดสอบ (Test Coverage)", [
        "## Authentication (9 cases)",
        "▸  Register ได้รับ JWT · ป้องกัน username/email ซ้ำ · Login ผิดรหัส = 401",
        "---",
        "## Projects (8 cases)",
        "▸  Worker สร้างโปรเจกต์ไม่ได้ (403) · Manager สร้างได้ · RBAC ถูกต้อง",
        "---",
        "## Tasks & Status Machine (6 cases)",
        "▸  สร้างงาน Open → In Progress บันทึก started_at อัตโนมัติ",
        "▸  เปลี่ยนเป็น Closed บันทึก completed_at · Worker ปิดงานเองไม่ได้",
        "---",
        "## Messaging (5 cases)",
        "▸  ส่งข้อความใน Project · ดึงประวัติ · บันทึก Read Receipt",
    ], ACCENT_TEAL, "📊")

    # ── 09  Profiling Intro ────────────────────────────────────────
    slide_section(prs, "Phase 4: Profiling — วัดก่อน แก้ทีหลัง", [
        '## หลักการ: "You can\'t improve what you don\'t measure"',
        "---",
        "## Static Profiling  (วิเคราะห์โค้ดก่อนรัน)",
        '▸  Next.js Production Build  →  ดู Bundle Size ต่อหน้า',
        '▸  npm audit  →  ตรวจหา Security Vulnerability ใน Dependencies',
        "---",
        "## Dynamic Profiling  (วิเคราะห์ขณะรัน)",
        '▸  Chrome DevTools — Flame Graph  →  ดู CPU usage & FPS',
        '▸  Heap Snapshot  →  ตรวจสอบ Memory Leak',
        '▸  API Response Time  →  วัดผ่าน Express Middleware Timer',
    ], ACCENT_INDIGO, "📈")

    # ── 10  Profiling Results Phase 3 vs 4 ────────────────────────
    slide_two_col(prs,
        "ผล Profiling: Phase 3 → Phase 4",
        [
            "## Phase 3 (ก่อน)",
            "▸ Shared Bundle: ~84.3 kB",
            "▸ Page (Project Detail): 13.4 kB",
            "▸ 1 Critical Vulnerability (SSRF)",
            "▸ API Response: ~50ms avg",
            "▸ FPS: ไม่ได้วัด",
            "▸ Memory Leak: ไม่ได้ตรวจสอบ",
        ],
        [
            "## Phase 4 (หลัง Optimize)",
            "▸ Shared Bundle: 87.4 kB (stable)",
            "▸ Page (Gantt): 6.66 kB  ✅ -50%",
            "▸ Vulnerability: Patched ✅",
            "▸ API Response: 15–30ms   ✅ -50%",
            "▸ FPS: 60 FPS stable     ✅",
            "▸ Memory Leak: ไม่พบ     ✅",
        ],
        ACCENT_INDIGO)

    # ── 11  Optimization Actions ───────────────────────────────────
    slide_section(prs, "การ Optimize จาก Profiling Result", [
        "## Static Fix",
        '▸  อัปเกรด Next.js  →  ปิดช่อง SSRF + Cache Poisoning',
        '▸  Dynamic Import (Code Splitting)  →  ลดขนาด Page Bundle',
        "---",
        "## Dynamic Fix",
        '▸  useMemo บน Gantt Calculation  →  ไม่ Re-compute ทุก Render',
        '▸  socket.off() ทุก Unmount  →  ป้องกัน Event Listener Leak',
        '▸  Async/Await ทุกจุด  →  ไม่มี Blocking Code บน Event Loop',
        "---",
        "▸  ทุกการแก้ไขมี Profiling Run ยืนยัน ก่อนและหลัง",
    ], ACCENT_INDIGO, "🔧")

    # ── 12  CI/CD Pipeline ─────────────────────────────────────────
    slide_section(prs, "CI/CD: Automate the Engineering Process", [
        "## Trigger: ทุกครั้งที่ git push → GitHub Actions เริ่มทำงานอัตโนมัติ",
        "---",
        "## Job 1 — Backend CI (Integration Testing)",
        "  checkout  →  npm ci  →  prisma generate  →  npm test",
        "  ✅ ถ้า 29/29 ผ่าน = โค้ดปลอดภัย",
        "---",
        "## Job 2 — Frontend CD Check (Build Validation)",
        "  checkout  →  npm ci  →  npm run build",
        "  ✅ ถ้า Build ผ่าน = Bundle Size ปกติ, ไม่มี Type Error",
        "---",
        "▸  Jobs ทั้งสองรัน Parallel  →  ประหยัดเวลา ~50%",
        "▸  ถ้าขั้นตอนไหนพัง  →  Block Merge ทันที (Zero Bug Policy)",
    ], ACCENT_BLUE, "🤖")

    # ── 13  CI/CD Value ────────────────────────────────────────────
    slide_section(prs, "ทำไม CI/CD ถึงสำคัญ?", [
        "## ปัญหาที่แก้ได้ด้วย CI/CD",
        '▸  "ในเครื่องผมใช้ได้" — แต่บน Server พัง  →  CI ตรวจสอบก่อน Merge',
        '▸  ลืมรัน Test ก่อน Push  →  Pipeline บังคับรันทุกครั้ง',
        '▸  Build Size พุ่งโดยไม่รู้  →  CD Job แจ้งเตือนทันที',
        "---",
        "## ในฐานะ SE: CI/CD = Safety Net",
        "▸  ทีมแก้โค้ดได้อย่างมั่นใจ เพราะมีหุ่นยนต์คอยตรวจ",
        "▸  ลดเวลา Manual QA ลง → โฟกัสที่ Feature ใหม่ได้เลย",
        "▸  บริษัท Tech ระดับโลก (Google, Meta) ใช้แนวทางนี้ทุกที่",
    ], ACCENT_BLUE, "💡")

    # ── 14  Known Issues & Backlog ─────────────────────────────────
    slide_section(prs, "Known Issues & Engineering Backlog", [
        "## ปัญหาที่พบและประเมินผลกระทบแล้ว",
        "▸  [Low]  Timezone Shift บน Gantt  →  แก้ด้วย date-fns UTC parser",
        "▸  [Low]  Dark Mode Hydration Flash  →  แก้ด้วย inline script บน <head>",
        "▸  [Med]  Socket Room ไม่ update แบบ Live  →  ใช้ emit('join-room') แก้ได้",
        "▸  [Med]  Storage Leak (Orphan Files)  →  เพิ่ม fs.unlink() บน Delete",
        "---",
        "## มุมมองวิศวกร",
        "▸  การระบุ Known Issue แสดงว่า 'รู้จักระบบตัวเองดี'",
        "▸  ทุก Issue มีแผน Fix ชัดเจน ไม่ใช่แค่พบแล้วทิ้งไว้",
        "▸  จัดลำดับความสำคัญ (Priority) ตาม Impact vs Effort",
    ], ACCENT_INDIGO, "⚠️")

    # ── 15  Summary ────────────────────────────────────────────────
    slide_section(prs, "สรุป: Engineering Excellence", [
        "## กระบวนการครบวงจร",
        "▸  Requirements → Design → Build → Test → Profile → Automate",
        "---",
        "## ตัวเลขที่พิสูจน์คุณภาพ",
        "▸  29/29 Tests Passed   |   83.6% Code Coverage",
        "▸  Bundle 6.66 kB/page  |   API Response 15–30ms",
        "▸  60 FPS UI            |   CI/CD Pipeline ✅ Green",
        "---",
        "## Shark Task ≠ แค่เว็บ",
        "▸  คือโปรเจกต์ที่ผ่านกระบวนการ Software Engineering จริงทุกขั้น",
        "▸  ทุกตัวเลขมาจากการวัดจริง ไม่ใช่การคาดเดา",
    ], ACCENT_BLUE, "🏆")

    prs.save("SharkTask_Final_Engineering.pptx")
    print("✅  Saved: SharkTask_Final_Engineering.pptx  (15 slides)")

if __name__ == "__main__":
    build()
