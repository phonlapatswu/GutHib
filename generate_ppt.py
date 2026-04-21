from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_slide_bg(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def set_text_color(tf, color_rgb):
    for paragraph in tf.paragraphs:
        paragraph.font.color.rgb = color_rgb

def create_engineering_presentation():
    prs = Presentation()

    # Themes
    shark_dark = RGBColor(15, 23, 42) # Slate 900
    shark_blue = RGBColor(59, 130, 246) # Blue 500
    shark_mint = RGBColor(94, 225, 205) # Mint
    white = RGBColor(255, 255, 255)
    gray = RGBColor(148, 163, 184)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SHARK TASK"
    title.text_frame.paragraphs[0].font.color.rgb = shark_mint
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "Engineering Deep-Dive: From Architecture to Optimization\nFinal Software Engineering Project Presentation"
    set_text_color(subtitle.text_frame, white)

    # --- Slide 2: SDLC Methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "1. SDLC: Iterative & Incremental Model"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "กระบวนการพัฒนาซอฟต์แวร์แบบก้าวกระโดด"
    p = tf.add_paragraph()
    p.text = "• Phase 1: Requirement Analysis - นิยามปัญหา 'GitHub for Everyone'"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Phase 2: Design & Prototyping - ออกแบบ Modular Architecture"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Phase 3: Iteration - พัฒนาทีละฟีเจอร์ (Auth -> Kanban -> Gantt)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Phase 4: Profiling & Verification - วัดผลและทำการ Optimization"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 3: System Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "2. Architecture Design Patterns"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "Separation of Concerns & Communication"
    p = tf.add_paragraph()
    p.text = "• Layered Pattern: แยก Controller (Logic) ออกจาก Model (Data)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real-time Layer: Socket.io Integration แยกส่วนจาก HTTP REST Pipeline"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Client-Side Logic: Axios Interceptor จัดการ Auth Token แบบ Centralized"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Theme Engine: ใช้ CSS Variables ควบคุม Dynamic UI (Dark/Light mode)"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 4: Data Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "3. Data Architecture & Relational Mapping"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "Relational Normalization with Prisma"
    p = tf.add_paragraph()
    p.text = "• Entities: User, Project, Task, Message, Comment, Log"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Junction Tables: ProjectsOnUsers (Many-to-Many สำหรับการแชร์โปรเจกต์)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Integrity Constraints: ใช้ Unique Indexes และ Foreign Key Cascading"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Database Indexing: กลยุทธ์การเพิ่ม B-Tree Index บนฟิลด์ Date เพื่อ Performance"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 5: Gantt Algorithm Deep-Dive ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "4. Technical Depth: Gantt Timeline Logic"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "Coordinate Systems in Web Development"
    p = tf.add_paragraph()
    p.text = "• Logic: คำนวณความต่างของมิลลิวินาที -> วัน -> แปลงเป็น CSS Grid X-Axis"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Formula: (TargetDate - ChartAnchor) / OneDayMS = Offset Column"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Optimization: พัฒนาระบบให้เป็น Modular Component แยกอิสระจาก Page Logic"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 6: QA & Verification Strategy ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "5. QA Strategy: The Multi-tier testing"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "ความมั่นใจผ่านกระบวนการตรวจสอบ (Coverage 83.6%)"
    p = tf.add_paragraph()
    p.text = "• Tier 1: Unit Test (Superset) - ตรวจสอบ Business Logic และ Data Schema"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Tier 2: API Testing (Supertest) - ตรวจสอบความถูกต้องของ HTTP Status & JSON"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Tier 3: UI Testing - ตรวจสอบ Workflow สำคัญ (Auth, Kanban Transition)"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 7: Automation & CI/CD ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "6. DevOps: Parallel CI/CD Pipeline"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "Efficiency in Development Workflow"
    p = tf.add_paragraph()
    p.text = "• Infrastructure: GitHub Actions Workflow"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Parallelization Job: รัน Backend Test และ Frontend Build พร้อมกัน"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Benefit: ลด Pipeline Time ลง 50% ทำให้การทำ Continuous Integration ไหลลื่น"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 8: The Profiling Methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "7. Profiling: Measure -> Analyze -> Fix"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "กระบวนการจูนระบบด้วยข้อมูลจริง"
    p = tf.add_paragraph()
    p.text = "• Static Audit: ใช้ npm audit ค้นหา Security Vulnerabilities"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Production Build Audit: วิเคราะห์ขนาด JS Bundle และความเร็วในการ Compile"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Runtime Analysis: ตรวจวัด Network Latency และการแสดงผล UI"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 9: Technical Optimizations (Deep-Dive) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "8. Optimization: Technical Rationales"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "ทำไมและอย่างไร? (The Engineering behind fixes)"
    p = tf.add_paragraph()
    p.text = "• Dynamic Code Splitting: ลด TTI (Time to Interactive) โดยโหลด Chart เฉพาะเมื่อเรียกใช้"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Security Hardening: อุดช่องโหว่ความรุนแรงระดับ Critical (SSRF/Cache Poisoning)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Next/Image Migration: ทำ Automatic Format Conversion (WebP) และ Resizing"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 10: Verification Metrics ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "9. Verification: Final Success Metrics"
    title.text_frame.paragraphs[0].font.color.rgb = shark_mint
    tf = slide.placeholders[1].text_frame
    tf.text = "การบรรลุเป้าหมายทางวิศวกรรม"
    p = tf.add_paragraph()
    p.text = "• Security: ปิดช่องโหว่ Critical สำเร็จ 100%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Performance: ลด JS Bundle ของหน้าโปรเจกต์ลงได้ 50% (จาก 13.4 -> 6.6 KB)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Build Speed: เร็วขึ้น 14% (ลงมาเหลือ 9.45s)"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 11: Risk Management (The Bug Log) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "10. Risk Management & Backlog"
    title.text_frame.paragraphs[0].font.color.rgb = shark_blue
    tf = slide.placeholders[1].text_frame
    tf.text = "การบริหารจัดการข้อผิดพลาดและแผนในอนาคต"
    p = tf.add_paragraph()
    p.text = "• Identified Risks: Socket Sync issues, Timezone Shifts, Storage Leak"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Mitigation Plan: ออกแบบ Cleanup Logic สำหรับไฟล์ และย้ายไปใช้ UTC บน DB"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Agile Backlog: จัดลำดับความสำคัญของ Bug รายงานให้ชัดเจนว่าเป็นหนทางพัฒนาต่อ"
    p.level = 1
    set_text_color(tf, white)

    # --- Slide 12: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, shark_dark)
    title = slide.shapes.title
    title.text = "11. SE Best Practices Summary"
    title.text_frame.paragraphs[0].font.color.rgb = shark_mint
    tf = slide.placeholders[1].text_frame
    tf.text = "บทสรุปหัวใจสำคัญของโปรเจกต์"
    p = tf.add_paragraph()
    p.text = "• SoC (Separation of Concerns): หัวใจของการดูแลรักษาระบบง่าย"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Data-Driven Optimization: ปรับปรุงระบบตามผลการ Profiling จริง"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Maintenance Matters: การทำ Documentation และรายงาน Bug อย่างตรงไปตรงมา"
    p.level = 1
    set_text_color(tf, white)

    prs.save('SharkTask_Engineering_Final.pptx')
    print("Final engineering presentation created successfully as SharkTask_Engineering_Final.pptx")

if __name__ == "__main__":
    create_engineering_presentation()
